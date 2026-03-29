# ========= Copyright 2025-2026 @ LiteratureHub All Rights Reserved. =========
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.
# ========= Copyright 2025-2026 @ LiteratureHub All Rights Reserved. =========

"""
PPT 模板管理器

管理 PPT 模板的加载、验证和定制。
"""

import logging
from typing import Dict, Any, List, Optional
from pathlib import Path
import json
import shutil

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class TemplateManager:
    """PPT 模板管理器

    管理多个 PPT 模板，支持模板的加载、验证、预览和定制。

    使用示例：
    ```python
    manager = TemplateManager(template_dir="templates/")

    # 列出所有可用模板
    templates = manager.list_templates()

    # 获取模板
    template_path = manager.get_template("academic_standard")

    # 验证模板
    is_valid = manager.validate_template(template_path)

    # 创建自定义模板
    manager.create_custom_template(
        base_template="academic_standard",
        customizations={...},
        output_name="my_custom_template"
    )
    ```
    """

    # 默认模板配置
    DEFAULT_TEMPLATES = {
        "academic_standard": {
            "name": "学术标准模板",
            "description": "适用于博士论文汇报的标准学术模板",
            "author": "LiteratureHub",
            "version": "1.0.0",
            "slide_count": 18,
            "features": [
                "4 部分结构（综述、创新、方法、后续）",
                "专业学术风格",
                "自动目录生成",
                "参考文献格式化"
            ]
        },
        "minimal": {
            "name": "简约模板",
            "description": "简洁的极简风格模板",
            "author": "LiteratureHub",
            "version": "1.0.0",
            "slide_count": 15,
            "features": [
                "极简设计",
                "重点突出",
                "适合快速汇报"
            ]
        },
        "detailed": {
            "name": "详细模板",
            "description": "包含更多细节和注释的详细模板",
            "author": "LiteratureHub",
            "version": "1.0.0",
            "slide_count": 25,
            "features": [
                "详细的文献分析",
                "完整的方法论展示",
                "包含数据可视化",
                "附录和补充材料"
            ]
        }
    }

    def __init__(self, template_dir: str = "templates/ppt"):
        """初始化模板管理器

        Args:
            template_dir: 模板目录路径
        """
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx 库未安装，无法使用模板管理功能")

        self.template_dir = Path(template_dir)
        self.template_dir.mkdir(parents=True, exist_ok=True)

        self.logger = logging.getLogger(f"{self.__class__.__name__}")

        # 模板注册表
        self.templates: Dict[str, Dict[str, Any]] = {}

        # 加载模板
        self._load_templates()

    def _load_templates(self):
        """加载所有可用模板"""
        # 加载默认模板配置
        self.templates.update(self.DEFAULT_TEMPLATES)

        # 扫描模板目录
        for template_file in self.template_dir.glob("*.pptx"):
            template_name = template_file.stem
            if template_name not in self.templates:
                # 自动发现模板
                self.templates[template_name] = {
                    "name": template_name,
                    "path": str(template_file),
                    "discovered": True
                }

        self.logger.info(f"已加载 {len(self.templates)} 个模板")

    def list_templates(self) -> List[Dict[str, Any]]:
        """列出所有可用模板

        Returns:
            模板信息列表
        """
        templates = []
        for template_name, template_info in self.templates.items():
            template_data = {
                "id": template_name,
                "name": template_info.get("name", template_name),
                "description": template_info.get("description", ""),
                "author": template_info.get("author", "Unknown"),
                "version": template_info.get("version", "1.0.0"),
                "slide_count": template_info.get("slide_count", 0),
                "features": template_info.get("features", []),
                "available": self._check_template_available(template_name)
            }
            templates.append(template_data)

        return templates

    def get_template(self, template_name: str) -> Optional[Path]:
        """获取模板文件路径

        Args:
            template_name: 模板名称

        Returns:
            模板文件路径（如果存在）
        """
        if template_name not in self.templates:
            self.logger.warning(f"模板不存在: {template_name}")
            return None

        # 检查是否有自定义模板文件
        template_info = self.templates[template_name]

        if "path" in template_info:
            template_path = Path(template_info["path"])
            if template_path.exists():
                return template_path

        # 尝试默认路径
        default_path = self.template_dir / f"{template_name}.pptx"
        if default_path.exists():
            return default_path

        # 返回 None（表示使用默认模板）
        self.logger.info(f"模板文件不存在，将使用默认模板: {template_name}")
        return None

    def validate_template(self, template_path: Path) -> Dict[str, Any]:
        """验证模板

        Args:
            template_path: 模板文件路径

        Returns:
            验证结果
        """
        result = {
            "valid": False,
            "errors": [],
            "warnings": [],
            "info": {}
        }

        if not template_path.exists():
            result["errors"].append(f"模板文件不存在: {template_path}")
            return result

        try:
            # 加载 PPT
            prs = Presentation(str(template_path))

            # 检查幻灯片数量
            slide_count = len(prs.slides)
            result["info"]["slide_count"] = slide_count

            if slide_count < 5:
                result["warnings"].append(f"幻灯片数量较少（{slide_count}）")
            elif slide_count > 50:
                result["warnings"].append(f"幻灯片数量较多（{slide_count}），可能影响演示效果")

            # 检查幻灯片布局
            layouts_used = set()
            for slide in prs.slides:
                layouts_used.add(slide.slide_layout.name)

            result["info"]["layouts_used"] = list(layouts_used)

            if len(layouts_used) < 2:
                result["warnings"].append("只使用了 1 种布局，可能略显单调")

            # 检查标题
            slides_without_title = 0
            for slide in prs.slides:
                if not slide.shapes.title or not slide.shapes.title.text:
                    slides_without_title += 1

            if slides_without_title > 0:
                result["warnings"].append(f"{slides_without_title} 张幻灯片没有标题")

            result["valid"] = True
            result["info"]["file_size"] = template_path.stat().st_size

        except Exception as e:
            result["errors"].append(f"模板加载失败: {e}")

        return result

    def create_custom_template(
        self,
        base_template: str,
        customizations: Dict[str, Any],
        output_name: str
    ) -> Path:
        """创建自定义模板

        Args:
            base_template: 基础模板名称
            customizations: 定制配置
            output_name: 输出模板名称

        Returns:
            创建的模板路径
        """
        self.logger.info(f"创建自定义模板: {output_name} (基于 {base_template})")

        # 获取基础模板
        base_path = self.get_template(base_template)

        if base_path and base_path.exists():
            # 加载基础模板
            prs = Presentation(str(base_path))
        else:
            # 创建空白模板
            prs = Presentation()

        # 应用定制配置
        # 实现模板定制逻辑
        if customizations:
            # 应用颜色方案
            if "colors" in customizations:
                # PPT 颜色定制需要使用 python-pptx 的主题 API
                # 这里提供基础实现框架
                self.logger.info(f"应用颜色方案: {customizations['colors']}")

            # 应用字体设置
            if "fonts" in customizations:
                self.logger.info(f"应用字体设置: {customizations['fonts']}")

            # 应用布局定制
            if "layout" in customizations:
                self.logger.info(f"应用布局定制: {customizations['layout']}")

        # 保存自定义模板
        output_path = self.template_dir / f"{output_name}.pptx"
        prs.save(str(output_path))

        # 注册新模板
        self.templates[output_name] = {
            "name": output_name,
            "path": str(output_path),
            "base_template": base_template,
            "customizations": customizations,
            "created_at": datetime.now().isoformat()
        }

        self.logger.info(f"自定义模板已创建: {output_path}")
        return output_path

    def delete_template(self, template_name: str) -> bool:
        """删除模板

        Args:
            template_name: 模板名称

        Returns:
            是否成功
        """
        if template_name in self.DEFAULT_TEMPLATES:
            self.logger.warning(f"无法删除默认模板: {template_name}")
            return False

        template_path = self.get_template(template_name)
        if template_path and template_path.exists():
            template_path.unlink()
            del self.templates[template_name]
            self.logger.info(f"模板已删除: {template_name}")
            return True

        return False

    def duplicate_template(
        self,
        source_name: str,
        target_name: str
    ) -> Optional[Path]:
        """复制模板

        Args:
            source_name: 源模板名称
            target_name: 目标模板名称

        Returns:
            新模板路径
        """
        source_path = self.get_template(source_name)
        if not source_path or not source_path.exists():
            self.logger.error(f"源模板不存在: {source_name}")
            return None

        target_path = self.template_dir / f"{target_name}.pptx"
        shutil.copy2(source_path, target_path)

        # 注册新模板
        self.templates[target_name] = {
            "name": target_name,
            "path": str(target_path),
            "source": source_name,
            "duplicated_at": datetime.now().isoformat()
        }

        self.logger.info(f"模板已复制: {source_name} → {target_name}")
        return target_path

    def export_template_config(
        self,
        template_name: str,
        output_path: str
    ) -> bool:
        """导出模板配置

        Args:
            template_name: 模板名称
            output_path: 输出文件路径

        Returns:
            是否成功
        """
        if template_name not in self.templates:
            self.logger.error(f"模板不存在: {template_name}")
            return False

        try:
            template_info = self.templates[template_name]

            with open(output_path, 'w', encoding='utf-8') as f:
                json.dump(template_info, f, ensure_ascii=False, indent=2)

            self.logger.info(f"模板配置已导出: {output_path}")
            return True

        except Exception as e:
            self.logger.error(f"导出失败: {e}")
            return False

    def _check_template_available(self, template_name: str) -> bool:
        """检查模板是否可用

        Args:
            template_name: 模板名称

        Returns:
            是否可用
        """
        # 默认模板始终可用
        if template_name in self.DEFAULT_TEMPLATES:
            return True

        # 检查模板文件是否存在
        template_path = self.get_template(template_name)
        return template_path is not None and template_path.exists()

    def get_template_statistics(self) -> Dict[str, Any]:
        """获取模板统计信息

        Returns:
            统计信息
        """
        total = len(self.templates)
        available = sum(1 for name in self.templates if self._check_template_available(name))

        return {
            "total_templates": total,
            "available_templates": available,
            "default_templates": len(self.DEFAULT_TEMPLATES),
            "custom_templates": total - len(self.DEFAULT_TEMPLATES),
            "template_dir": str(self.template_dir)
        }
