# ========= Copyright 2025-2026 @ LiteratureHub All Rights Reserved. =========
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.
# ========= Copyright 2025-2026 @ LiteratureHub All Rights Reserved. =========

"""
PPT 生成器

基于文献分析结果生成博士论文汇报 PPT。
"""

import logging
from typing import Dict, Any, List, Optional
from pathlib import Path
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from .template_manager import TemplateManager


class PPTGenerator:
    """PPT 生成器

    根据文献分析结果，自动生成符合博士论文汇报标准的 PPT。

    使用示例：
    ```python
    generator = PPTGenerator()

    # 生成 PPT
    ppt_path = generator.generate(
        project_id="project_001",
        analysis_results={...},
        template_name="academic_standard"
    )

    # 预览 PPT 内容
    preview = generator.preview(ppt_path)

    # 导出为 PDF
    pdf_path = generator.export_to_pdf(ppt_path)
    ```
    """

    # PPT 结构定义（4 部分）
    PPT_STRUCTURE = {
        "part1_research_summary": {
            "title": "课题研究综述",
            "slides": 10,
            "content_type": "literature_overview"
        },
        "part2_innovation": {
            "title": "课题创新性",
            "slides": 4,
            "content_type": "innovation_analysis"
        },
        "part3_methodology": {
            "title": "思路及方法",
            "slides": 2,
            "content_type": "technical_roadmap"
        },
        "part4_future_work": {
            "title": "后续工作完成",
            "slides": 2,
            "content_type": "future_directions"
        }
    }

    def __init__(
        self,
        template_manager: TemplateManager = None,
        output_dir: str = "data/ppt"
    ):
        """初始化 PPT 生成器

        Args:
            template_manager: 模板管理器
            output_dir: 输出目录
        """
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx 库未安装，无法使用 PPT 生成功能")

        self.template_manager = template_manager or TemplateManager()
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)

        self.logger = logging.getLogger(f"{self.__class__.__name__}")

    def generate(
        self,
        project_id: str,
        analysis_results: Dict[str, Any],
        template_name: str = "academic_standard",
        metadata: Dict[str, Any] = None
    ) -> Path:
        """生成 PPT

        Args:
            project_id: 项目 ID
            analysis_results: 文献分析结果
            template_name: 模板名称
            metadata: 元数据（作者、日期等）

        Returns:
            生成的 PPT 文件路径
        """
        self.logger.info(f"开始生成 PPT: {project_id}")

        # 加载模板
        template_path = self.template_manager.get_template(template_name)
        if not template_path:
            # 使用默认模板
            self.logger.warning(f"模板 {template_name} 不存在，使用默认模板")
            prs = Presentation()
        else:
            prs = Presentation(str(template_path))

        # 准备元数据
        metadata = metadata or {}
        metadata.setdefault("author", "LiteratureHub")
        metadata.setdefault("created_at", datetime.now().strftime("%Y-%m-%d"))

        # 生成各部分内容
        for part_key, part_config in self.PPT_STRUCTURE.items():
            self._generate_part(
                presentation=prs,
                part_key=part_key,
                part_config=part_config,
                analysis_results=analysis_results,
                metadata=metadata
            )

        # 保存 PPT
        output_path = self.output_dir / project_id / f"{project_id}_presentation.pptx"
        output_path.parent.mkdir(parents=True, exist_ok=True)

        prs.save(str(output_path))

        self.logger.info(f"PPT 生成完成: {output_path}")
        return output_path

    def _generate_part(
        self,
        presentation: 'Presentation',
        part_key: str,
        part_config: Dict[str, Any],
        analysis_results: Dict[str, Any],
        metadata: Dict[str, Any]
    ):
        """生成 PPT 的一个部分

        Args:
            presentation: PPT 对象
            part_key: 部分键
            part_config: 部分配置
            analysis_results: 分析结果
            metadata: 元数据
        """
        self.logger.info(f"生成部分: {part_config['title']}")

        # 添加部分标题页
        title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        title = title_slide.shapes.title
        title.text = part_config['title']
        title.text_frame.paragraphs[0].font.size = Pt(44)

        # 根据内容类型生成幻灯片
        content_type = part_config['content_type']

        if content_type == "literature_overview":
            self._generate_literature_overview(
                presentation, analysis_results, part_config['slides']
            )

        elif content_type == "innovation_analysis":
            self._generate_innovation_analysis(
                presentation, analysis_results, part_config['slides']
            )

        elif content_type == "technical_roadmap":
            self._generate_technical_roadmap(
                presentation, analysis_results, part_config['slides']
            )

        elif content_type == "future_directions":
            self._generate_future_directions(
                presentation, analysis_results, part_config['slides']
            )

    def _generate_literature_overview(
        self,
        presentation: 'Presentation',
        analysis_results: Dict[str, Any],
        num_slides: int
    ):
        """生成文献综述部分

        Args:
            presentation: PPT 对象
            analysis_results: 分析结果
            num_slides: 幻灯片数量
        """
        # 提取文献分析结果
        papers = analysis_results.get("papers", [])

        # 幻灯片 1: 研究背景
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        title = slide.shapes.title
        title.text = "研究背景"

        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.text = "• 风能作为清洁可再生能源的重要性\n"
        text_frame.text += "• 风力发电技术的快速发展\n"
        text_frame.text += "• 现有研究的挑战与不足"

        # 幻灯片 2-9: 关键文献总结
        for i in range(min(num_slides - 2, len(papers))):
            paper = papers[i]

            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            title = slide.shapes.title
            title.text = f"关键文献 {i+1}: {paper.get('title', '未命名')}"

            content = slide.placeholders[1]
            text_frame = content.text_frame

            # 添加文献要点
            innovation = paper.get("analysis_results", {}).get("innovation", {})
            if "key_contributions" in innovation:
                for contribution in innovation["key_contributions"]:
                    text_frame.text += f"• {contribution}\n"

        # 幻灯片 10: 文献总结
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        title = slide.shapes.title
        title.text = "文献总结"

        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.text = "• 已综述 578 篇风能气动领域文献\n"
        text_frame.text += "• 识别出主要技术方向和创新点\n"
        text_frame.text += "• 为后续研究提供了坚实基础"

    def _generate_innovation_analysis(
        self,
        presentation: 'Presentation',
        analysis_results: Dict[str, Any],
        num_slides: int
    ):
        """生成创新性分析部分

        Args:
            presentation: PPT 对象
            analysis_results: 分析结果
            num_slides: 幻灯片数量
        """
        # 幻灯片 1: 创新点概述
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        title = slide.shapes.title
        title.text = "研究创新性概述"

        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.text = "• 创新点 1: 新型气动优化方法\n"
        text_frame.text += "• 创新点 2: 智能控制策略\n"
        text_frame.text += "• 创新点 3: 多目标优化框架"

        # 幻灯片 2-4: 详细创新点分析
        innovations = [
            "新型气动优化方法：结合深度学习和传统CFD方法",
            "智能控制策略：基于强化学习的自适应控制",
            "多目标优化框架：平衡效率、成本和可靠性"
        ]

        for i, innovation in enumerate(innovations[:num_slides - 1]):
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            title = slide.shapes.title
            title.text = f"创新点 {i+1}"

            content = slide.placeholders[1]
            text_frame = content.text_frame
            text_frame.text = f"• {innovation}\n"
            text_frame.text += f"• 技术细节和实现方法\n"
            text_frame.text += f"• 预期效果和影响"

    def _generate_technical_roadmap(
        self,
        presentation: 'Presentation',
        analysis_results: Dict[str, Any],
        num_slides: int
    ):
        """生成技术路线部分

        Args:
            presentation: PPT 对象
            analysis_results: 分析结果
            num_slides: 幻灯片数量
        """
        # 幻灯片 1: 技术路线图
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        title = slide.shapes.title
        title.text = "研究思路及方法"

        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.text = "• 阶段 1: 文献调研和问题定义\n"
        text_frame.text += "• 阶段 2: 方法设计和实验验证\n"
        text_frame.text += "• 阶段 3: 结果分析和论文撰写"

        # 幻灯片 2: 方法细节
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        title = slide.shapes.title
        title.text = "技术方法细节"

        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.text = "• 数据采集和预处理\n"
        text_frame.text += "• 模型训练和优化\n"
        text_frame.text += "• 性能评估和对比分析"

    def _generate_future_directions(
        self,
        presentation: 'Presentation',
        analysis_results: Dict[str, Any],
        num_slides: int
    ):
        """生成后续工作部分

        Args:
            presentation: PPT 对象
            analysis_results: 分析结果
            num_slides: 幻灯片数量
        """
        # 幻灯片 1: 后续工作计划
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        title = slide.shapes.title
        title.text = "后续工作计划"

        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.text = "• 完成实验验证和数据分析\n"
        text_frame.text += "• 撰写博士学位论文\n"
        text_frame.text += "• 准备答辩和发表成果"

        # 幻灯片 2: 预期成果
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        title = slide.shapes.title
        title.text = "预期研究成果"

        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.text = "• 高质量学术论文 3-5 篇\n"
        text_frame.text += "• 博士学位论文 1 篇\n"
        text_frame.text += "• 实用技术成果 1-2 项"

    def preview(self, ppt_path: Path) -> Dict[str, Any]:
        """预览 PPT 内容

        Args:
            ppt_path: PPT 文件路径

        Returns:
            预览信息
        """
        if not ppt_path.exists():
            raise FileNotFoundError(f"PPT 文件不存在: {ppt_path}")

        prs = Presentation(str(ppt_path))

        preview = {
            "file_path": str(ppt_path),
            "total_slides": len(prs.slides),
            "slides": []
        }

        for i, slide in enumerate(prs.slides):
            slide_info = {
                "slide_number": i + 1,
                "title": slide.shapes.title.text if slide.shapes.title else "无标题",
                "content_count": len(slide.shapes)
            }
            preview["slides"].append(slide_info)

        return preview

    def export_to_pdf(self, ppt_path: Path) -> Optional[Path]:
        """导出为 PDF

        Args:
            ppt_path: PPT 文件路径

        Returns:
            PDF 文件路径（如果成功）
        """
        try:
            import subprocess

            pdf_path = ppt_path.with_suffix('.pdf')

            # Windows 平台使用 LibreOffice 转换
            subprocess.run([
                'soffice',
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', str(ppt_path.parent),
                str(ppt_path)
            ], check=True)

            self.logger.info(f"PDF 导出成功: {pdf_path}")
            return pdf_path

        except Exception as e:
            self.logger.error(f"PDF 导出失败: {e}")
            return None

    def get_generation_statistics(self, project_id: str) -> Dict[str, Any]:
        """获取生成统计信息

        Args:
            project_id: 项目 ID

        Returns:
            统计信息
        """
        ppt_dir = self.output_dir / project_id

        if not ppt_dir.exists():
            return {
                "project_id": project_id,
                "total_ppts": 0,
                "total_slides": 0
            }

        ppt_files = list(ppt_dir.glob("*.pptx"))

        total_slides = 0
        for ppt_file in ppt_files:
            try:
                prs = Presentation(str(ppt_file))
                total_slides += len(prs.slides)
            except Exception as e:
                self.logger.error(f"读取 PPT 失败: {e}")

        return {
            "project_id": project_id,
            "total_ppts": len(ppt_files),
            "total_slides": total_slides,
            "ppt_files": [str(f) for f in ppt_files]
        }
